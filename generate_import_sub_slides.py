#!/usr/bin/env python3
"""
Generate PowerPoint slides: India's Import Substitution Roadmap
Comprehensive presentation with 18 slides covering all scenarios and comparisons
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 45, 90)  # Dark blue

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)

    return slide

def add_content_slide(prs, title, content_lines, bg_color=RGBColor(245, 245, 245)):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 45, 90)

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(content_lines):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """Add two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 45, 90)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(6))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 90, 45)

    for line in left_content:
        p = left_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(3)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(6))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 30, 30)

    for line in right_content:
        p = right_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_before = Pt(3)

    return slide

# ── SLIDE 1: Title ──
add_title_slide(prs,
    "India's Import Substitution Roadmap",
    "FY26-FY30: Three Pathways to ₹75,000–95,000 Crores Annual Savings")

# ── SLIDE 2: Executive Summary ──
add_content_slide(prs, "Executive Summary: Quantum of Savings",
[
    "🎯 Realistic Base Case (70% execution): ₹75,000–90,000 cr/year ($9.0–10.8bn/year)",
    "",
    "🚀 Optimistic (All capex on-time): ₹98,000–130,000 cr/year ($11.8–15.7bn/year)",
    "",
    "🛡️ Conservative (1-2 delays): ₹50,000–65,000 cr/year ($6.0–7.8bn/year)",
    "",
    "Key insight: Real, structural import reductions + new export earnings, not projections"
])

# ── SLIDE 3: Three Pathways Overview ──
add_content_slide(prs, "Three Interconnected Pathways",
[
    "1️⃣ PETROCHEMICAL FEEDSTOCKS (Textiles) ────────── ₹50,000–95,000 cr/year",
    "   $20.4bn annual textile feedstock imports → 40-60% substitution via capex",
    "",
    "2️⃣ ETHANOL BLENDING & PETROL EXPORT ───────────── ₹12,000–18,000 cr/year",
    "   E20→E30 ramp frees domestic petrol for export ($7.3bn/year)",
    "",
    "3️⃣ RENEWABLE FUELS & GAS SUBSTITUTION ────────── ₹10,000–14,000 cr/year",
    "   CBG/SATAT scales 5 MMTPA by FY30 (displaces $1.2–1.8bn LNG/year)"
])

# ── SLIDE 4: Pathway 1 - Problem ──
add_content_slide(prs, "Pathway 1: Textile Feedstock Crisis",
[
    "💔 Current Import Dependency (FY25):",
    "  • Polymer imports (Ch 39): $22.1 billion/year (↑10% YoY)",
    "  • Chemical imports (Ch 29): $26.6 billion/year (↑10.8% YoY)",
    "  • Total feedstock gap: $48.7 billion/year (GROWING)",
    "",
    "🔴 Impact: Textile margins compressed, export competitiveness eroding",
    "  • INDUTECH: -80.9% margins (Saudi SABIC LLDPE @ $950-1000/MT)"
])

# ── SLIDE 5: Pathway 1 - Solution ──
add_content_slide(prs, "Pathway 1: Capex Solution ($38.5bn)",
[
    "✅ Government + Private Capex by FY30:",
    "",
    "  🏭 BPCL Andhra Pradesh: ₹1,00,000 cr ($11.4bn) → 1.8 MMTPA (FY27-30)",
    "  🏭 RIL Oil-to-Chemicals: ₹75,000 cr ($8.6bn) → 2.0 MMTPA (FY28-30)",
    "  🏭 L&T BPCL Bina: ₹5-10k cr ($600-1.2bn) → 1.0 MMTPA (FY28-29)",
    "  🏭 BPCL Kochi PP: ₹5,514 cr ($599mn) → 0.15 MMTPA (Commissioned 2026)",
    "  🏭 BHAVYA Parks: ₹3,030 cr ($365mn) → 3 sites + specialty chemicals",
    "",
    "Result: Polypropylene import dependency ELIMINATED; Polyethylene 70-80% closed"
])

# ── SLIDE 6: Pathway 1 - Segment Savings ──
add_content_slide(prs, "Pathway 1: Segment-Level Savings by FY30",
[
    "📊 INDUTECH:  ₹19,920 cr/year ($2.4bn) — BPCL Bina LLDPE @ $150-200/MT",
    "",
    "📊 HOMETECH: ₹5,478 cr/year ($0.66bn) — PE/PET cost reduction",
    "",
    "📊 MOBILTECH: ₹5,395 cr/year ($0.65bn) — Specialty polymers",
    "",
    "📊 BUILDTECH: ₹3,984 cr/year ($0.48bn) — Geotextiles",
    "",
    "📊 PACKTECH: ₹2,241 cr/year ($0.27bn) — Export-led growth",
    "",
    "📊 TOTAL: ₹50,298–95,450 cr/year (Base to Optimistic)"
])

# ── SLIDE 7: Pathway 2 - Ethanol ──
add_content_slide(prs, "Pathway 2: Ethanol Blending & Export Value",
[
    "🌾 Mechanism: Ethanol blends into petrol pool → vehicles buy more litres",
    "",
    "📈 FY30 E27 Scenario (Recommended):",
    "  • Freed petrol export value: ₹62,300 cr ($7.3bn/year) ← HEADLINE",
    "  • OMC pump margin on extra volume: ₹1,000 cr",
    "  • Tax revenue (excise + VAT): ₹11,500 cr",
    "  • State SGST (5%, price-neutral): ₹4,787 cr ← NEW POLICY LEVER",
    "",
    "💰 Total E27 pathway: ₹79,587 cr/year ($9.6bn/year at FY30)"
])

# ── SLIDE 8: Pathway 3 - Renewables ──
add_content_slide(prs, "Pathway 3: Renewable Fuels & CBG",
[
    "🌱 CBG (Compressed Biogas) via SATAT Programme:",
    "  • Current: 0.5 MMTPA (12-15 plants operating)",
    "  • FY30 Target: 5 MMTPA (5,000 plants planned)",
    "",
    "💨 LNG Import Displacement:",
    "  • Current LNG imports: ~180 MMTPA equivalent (~₹80,000 cr/year cost)",
    "  • CBG substitution by FY30: 10-15% (18 MMTPA)",
    "  • Annual savings: ₹8,000–12,000 crores ($1.0–1.4bn)",
    "",
    "🌾 Secondary: Bio-isobutanol (SAF route) ₹2-3k cr/year + Distillery conversion"
])

# ── SLIDE 9: Consolidated Savings ──
add_content_slide(prs, "Consolidated: All Pathways by FY30",
[
    "┌────────────────────────────────────────────────────────┐",
    "│ Petrochemical Feedstocks:      ₹50,000–95,000 cr/year  │",
    "│ Ethanol & Petrol Export:       ₹12,000–18,000 cr/year  │",
    "│ Renewable Fuels & CBG:         ₹10,000–14,000 cr/year  │",
    "├────────────────────────────────────────────────────────┤",
    "│ TOTAL ANNUAL (FY30):           ₹72,000–127,000 cr/year │",
    "│                                $8.7–15.3 billion/year  │",
    "├────────────────────────────────────────────────────────┤",
    "│ REALISTIC MID-POINT:           ₹75,000–95,000 cr/year  │",
    "│                                $9.0–11.5 billion/year  │",
    "└────────────────────────────────────────────────────────┘"
])

# ── SLIDE 10: Government Capex Required ──
add_content_slide(prs, "Investment Required: ₹1,38,514 Crores (~$16.6bn)",
[
    "Total Government + Private Capex (FY26-FY30):",
    "",
    "  ▪ BPCL Andhra Pradesh:         ₹1,00,000 cr (greenfield 1.8 MMTPA)",
    "  ▪ RIL Oil-to-Chemicals:        ₹75,000 cr (2.0 MMTPA integrated)",
    "  ▪ L&T BPCL Bina:               ₹5-10k cr (1.0 MMTPA swing)",
    "  ▪ BPCL Kochi PP:               ₹5,514 cr (0.15 MMTPA, commissioned)",
    "  ▪ BHAVYA Chemical Parks:       ₹3,030 cr (3 parks)",
    "  ▪ Private Sector:              ₹2-5k cr",
    "",
    "Expected Payback: 4-6 years post-FY30 | ROI: 15-25%"
])

# ── SLIDE 11: Current Import Bill Context ──
add_two_column_slide(prs, "Savings vs Current Baselines",
    "Current (FY25/CY2024)",
    [
        "Total imports: $718.16bn",
        "Trade deficit: $275.45bn",
        "Energy imports: $225.39bn",
        "Oil bill: $160.8bn",
        "Oil net bill: $116.4bn",
        "Textile feedstock: $48.7bn",
        "Forex addition: $10-20bn/yr"
    ],
    "Import Sub. Savings (FY30)",
    [
        "$9.0–11.5bn/year",
        "$9.0–11.5bn/year",
        "$8.5–9.1bn/year",
        "Partial (feedstock focus)",
        "$7.3–8.5bn/year",
        "$7.8–9.2bn/year",
        "$9.0–11.5bn/year"
    ])

# ── SLIDE 12: Impact %s ──
add_content_slide(prs, "Savings as % of Current Baselines: The Real Story",
[
    "🔴 TEXTILE FEEDSTOCK: 16–19% reduction ($7.8–9.2bn saved of $48.7bn)",
    "   → Most impactful sector; polypropylene import dependency ELIMINATED",
    "",
    "🔴 NET OIL BILL: 6.3–7.8% reduction ($7.3–8.5bn of $116.4bn)",
    "   → Most visible oil-sector leverage",
    "",
    "🟡 TOTAL TRADE DEFICIT: 3.3–4.2% improvement ($9.0–11.5bn of $275.45bn)",
    "   → Modest as % of deficit, but real structural reduction",
    "",
    "🟡 TOTAL IMPORTS: 1.3–1.6% reduction ($9.0–11.5bn of $718.16bn)",
    "   → Small percentage, large absolute value"
])

# ── SLIDE 13: Forex Impact ──
add_content_slide(prs, "Forex Reserves: The Strategic Multiplier",
[
    "Current Forex Reserves: ~$645–650bn (world's 4th largest)",
    "Typical Annual Addition: $10–20bn/year",
    "",
    "Import Substitution Contribution (FY30):",
    "  • Annual forex contribution: $9.0–11.5bn/year",
    "  • As % of typical annual addition: 45–115% ← EQUALS 1-2 years of normal growth",
    "",
    "5-Year Cumulative (FY26-FY30): $30–36 billion",
    "  • As % of typical 5-year addition ($50-100bn): 30–60%",
    "",
    "✅ Makes forex generation MORE RESILIENT (less dependent on volatile FDI/portfolio)"
])

# ── SLIDE 14: Critical Path Timeline ──
add_content_slide(prs, "Critical Path: FY26-FY30 Execution",
[
    "🔴 Q3-FY27: BPCL Andhra Pradesh cracker commissioning (load-bearing)",
    "",
    "🔴 FY28-29: L&T Bina + RIL O2C ramp-up (40-50% of savings hinge here)",
    "",
    "🟠 FY26-27: BHAVYA Parks site finalization + State SGST policy coordination",
    "",
    "🟠 FY26-FY30: SATAT CBG deployment (500+ plants needed for full benefit)",
    "",
    "🟡 FY29-30: Full run-rate capacity online; import duty escalation (10%→20%)"
])

# ── SLIDE 15: Key Risks ──
add_content_slide(prs, "Top 5 Execution Risks & Mitigation",
[
    "1️⃣ Refinery capex delays (BPCL/RIL) → Allocate ₹5k cr risk capital; fast-track approvals",
    "",
    "2️⃣ Crude < $60/bbl (uneconomical) → Hedging instruments; floor pricing guardrails",
    "",
    "3️⃣ Ethanol feedstock scarcity → Diversify (maize→FCI rice→molasses→stalk)",
    "",
    "4️⃣ China polyester dumping → Anti-dumping duties if domestic capacity idles",
    "",
    "5️⃣ State SGST coordination → Price-neutral design (₹4,787 cr new state revenue)"
])

# ── SLIDE 16: Policy Anchors ──
add_content_slide(prs, "Policy Enablers for FY26-FY30",
[
    "✅ IMMEDIATE (FY26-27):",
    "   • BHAVYA Parks site finalization + ₹500 cr release",
    "   • BPCL AP environmental clearance + construction start",
    "   • State SGST @ 5% on ethanol (price-neutral model)",
    "",
    "✅ MID-TERM (FY27-29):",
    "   • BPCL Kochi PP commissioning (trials → production)",
    "   • L&T Bina full construction + supplier contracts",
    "   • RIL O2C construction launch",
    "",
    "✅ FINAL (FY29-30):",
    "   • Import duty escalation (10%→20% on polymers, gradual)"
])

# ── SLIDE 17: Final Verdict ──
add_content_slide(prs, "Final Verdict: Structural Stabilizer",
[
    "✅ IS: A multi-year, structural import reduction + export enabler",
    "   ✓ Saves $9–11.5bn/year by FY30 in real import + export terms",
    "   ✓ Concentrates in textiles (16-19% feedstock reduction)",
    "   ✓ Builds forex resilience (less dependent on volatile FDI)",
    "",
    "❌ IS NOT: A silver bullet for trade deficit or forex crisis",
    "   ✗ Only 3-4% improvement in $275bn trade deficit",
    "   ✗ Only 1-3% of total imports; requires other levers too",
    "",
    "🎯 BEST USE: Pair with merchandise export growth + service sector + FDI + remittances"
])

# ── SLIDE 18: Call to Action ──
add_content_slide(prs, "Next Steps: Execution & Monitoring",
[
    "1. Lock in capex: BPCL AP, RIL O2C, L&T Bina (non-negotiable)",
    "",
    "2. Fast-track policy: State SGST design, import duties, anti-dumping framework",
    "",
    "3. Coordinate across ministries: Textiles, DPIIT, Petroleum, Finance",
    "",
    "4. Track quarterly: Capex progress, commissioning, import reductions, forex additions",
    "",
    "5. Adjust for risks: Capex delays, crude prices, feedstock availability, global demand"
], bg_color=RGBColor(200, 230, 255))

# Save presentation
output_path = "/Users/umashankar/IMPORT_SUBSTITUTION_ROADMAP.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
